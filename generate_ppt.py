from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, content_lines, is_title_slide=False):
    if is_title_slide:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = "\n".join(content_lines)
    else:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        tf.clear()
        
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line.replace("- ", "").replace("**", "")
            p.font.size = Pt(18)
            p.level = 0

prs = Presentation()

# Slide 1
add_slide(prs, "Nexus Workspace (Forge India Connect)", [
    "The Unified Real-Time Collaboration & Communication Platform",
    "Speaker: Your Name / Team Name",
    "Date: "
], True)

# Slide 2
add_slide(prs, "What is Nexus Workspace?", [
    "An all-in-one unified workspace platform designed to streamline team communication and collaboration.",
    "Bridges the gap between remote and hybrid teams with seamless real-time interactions.",
    "Combines video conferencing, instant messaging, email, and document collaboration into a single, cohesive ecosystem.",
    "Designed with both a comprehensive Web Application and a cross-platform Mobile Application."
])

# Slide 3
add_slide(prs, "A Multi-Platform Experience", [
    "Web Application:",
    "  Browser-based experience built on React and Vite.",
    "  Utilizes a Microfrontend Architecture, allowing applications like Chat, Meet, and Mail to operate independently yet cohesively.",
    "Mobile Application:",
    "  Native iOS and Android applications built with React Native and Expo.",
    "  Ensures teams stay connected on-the-go with optimized mobile layouts and hardware integrations."
])

# Slide 4
add_slide(prs, "Meet: High-Performance Video Meetings", [
    "Real-Time Media: Low-latency HD video and audio communication powered by WebRTC.",
    "Collaboration Tools: Seamless screen sharing capabilities across web and mobile.",
    "Persistent Rooms: Dedicated, always-on meeting rooms for continuous team collaboration.",
    "Host Controls: Advanced meeting moderation, participant muting, and room security.",
    "Smart Assistant: Integrated AI bot ('Forge India Connect AI') for real-time meeting transcription and post-meeting summaries."
])

# Slide 5
add_slide(prs, "Meet: How It Works", [
    "1. Lobby & Pre-flight: User schedules or starts a meeting, entering a lobby to configure camera and microphone settings.",
    "2. Signaling Phase: Client connects to the custom WebSocket server to exchange Session Descriptions (SDP) and ICE candidates.",
    "3. P2P Mesh Network: WebRTC establishes direct, secure, peer-to-peer media channels.",
    "4. Active Meeting: Users interact via video, voice, and in-meeting chat.",
    "5. Post-Meeting: The session ends, and the AI Smart Assistant processes the recorded audio to generate an automated meeting summary."
])

# Slide 6
add_slide(prs, "Chat: Instant Team Communication", [
    "Direct & Group Messaging: 1-on-1 private messages and dedicated group channels.",
    "Organized Threads: Reply specifically to messages to keep conversations organized and focused.",
    "Threads AI Assistant: Leverage integrated AI to automatically summarize long conversation threads, extract key action items, and suggest context-aware replies.",
    "Rich Media Sharing: Support for file attachments, images, and rich text formatting.",
    "Presence Indicators: Real-time online/offline status and typing indicators."
])

# Slide 7
add_slide(prs, "Chat: How It Works", [
    "1. Connection: Upon login, the client establishes a persistent WebSocket connection to the backend.",
    "2. Data Synchronization: The app fetches chat history from the Node.js/Fastify backend and stores it in the local state.",
    "3. Real-Time Delivery: When a user sends a message, it is instantly routed through the WebSocket server to the recipient's active session.",
    "4. AI Processing: For Threads AI actions, the backend compiles the specific thread context, securely queries the integrated LLM, and streams the summary back.",
    "5. Persistence: The backend asynchronously commits messages and AI summaries to the database, ensuring no data loss."
])

# Slide 8
add_slide(prs, "Mail: Integrated Email Client", [
    "Unified Inbox: Manage all internal and external communications directly within the workspace.",
    "Standard Mail Features: Inbox, Sent, Drafts, and Trash folders.",
    "Rich HTML Composer: Create beautifully formatted emails.",
    "AI-Assisted Drafting: Utilize AI to help compose, summarize, and reply to long email threads."
])

# Slide 9
add_slide(prs, "Mail: How It Works", [
    "1. Retrieval: The Mail microfrontend communicates with the backend via secure REST APIs to fetch emails via SMTP/IMAP protocols.",
    "2. Composition: Users draft emails using the integrated rich-text editor.",
    "3. AI Integration: If requested, the prompt is sent to the backend AI service which returns a drafted response.",
    "4. Dispatch: The email is sent through the secure backend routing system to external servers or internal users."
])

# Slide 10
add_slide(prs, "Collaborative Productivity Suite", [
    "Live Document Editing (Docs): Real-time collaborative text editing with rich formatting.",
    "Spreadsheet Manipulation (Sheets): Data organization, formulas, and real-time collaboration.",
    "Cloud Sync: Auto-saving to ensure work is never lost.",
    "Access Control: Share documents with specific read/write permissions."
])

# Slide 11
add_slide(prs, "Docs & Sheets: How It Works", [
    "1. Document Creation: A user creates a new file, which is initialized in the database.",
    "2. Live Synchronization: Using Operational Transformation (OT) or WebSockets, every keystroke or cell change is broadcast to all active viewers.",
    "3. Conflict Resolution: The backend acts as the source of truth, resolving simultaneous edits to prevent data corruption.",
    "4. Auto-Save: Periodic snapshots are saved to the persistent storage."
])

# Slide 12
add_slide(prs, "Super Admin & Team Management", [
    "Super Admin Dashboard: Comprehensive oversight of the entire workspace, user analytics, and system health.",
    "Team Management: Tools to onboard users, manage workspace permissions, and organize departments.",
    "Role-Based Access Control (RBAC): Strict permission models distinguishing between Hosts, Members, and Admins.",
    "Security & MFA: Management of Multi-Factor Authentication (MFA) and End-to-End Encryption (E2EE) policies."
])

# Slide 13
add_slide(prs, "Administration: How It Works", [
    "1. Secure Access: Admins log in using enhanced security measures (MFA).",
    "2. Analytics Aggregation: The backend aggregates data across all microfrontends (active meetings, total messages, active users) and serves it to the Dashboard.",
    "3. Policy Enforcement: When an Admin updates a user's role, the backend immediately invalidates old tokens and pushes a WebSocket event to enforce the new permissions on the client side in real-time."
])

# Slide 14
add_slide(prs, "The Technology Stack", [
    "Frontend: React.js, Vite, Tailwind CSS, Microfrontends architecture.",
    "Mobile: React Native, Expo, react-native-webrtc.",
    "Backend: Node.js, Fastify/Express, RESTful APIs.",
    "Real-Time / Media: WebSockets, WebRTC, STUN/TURN servers.",
    "AI & Security: Integrated AI models for transcription/summaries, robust JWT and MFA authentication."
])

# Slide 15
add_slide(prs, "Thank You!", [
    "Any Questions?",
    "Contact Information / Repository Links"
])

prs.save("Nexus_Workspace_Presentation.pptx")
print("Presentation generated successfully at Nexus_Workspace_Presentation.pptx")
